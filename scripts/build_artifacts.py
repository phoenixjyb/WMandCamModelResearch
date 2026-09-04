#!/usr/bin/env python3
"""Build reproducible reading artifacts from the canonical Markdown sources.

Outputs:
- DOCX and PDF editions for each research report / white paper
- editable PPTX decks from the bilingual slide Markdown

The renderer is intentionally dependency-light and deterministic. Markdown remains
canonical; derived files are convenience artifacts, not independent source text.
"""

from __future__ import annotations

import html
import re
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Iterator, Sequence

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Inches as PptxInches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    KeepTogether,
    PageBreak,
    PageTemplate,
    Paragraph,
    Preformatted,
    Spacer,
    Table,
    TableStyle,
)

ROOT = Path(__file__).resolve().parents[1]

WHITE_PAPERS = [
    ROOT / "docs/spatial-world-model-white-paper/v1.0/en/From_Camera_Controlled_Video_to_Spatial_World_Models_Recomo_White_Paper_v1.0_EN.md",
    ROOT / "docs/spatial-world-model-white-paper/v1.0/zh-CN/从相机可控视频到空间世界模型_Recomo技术白皮书_v1.0_CN.md",
    ROOT / "docs/spatial-world-model-white-paper/v2.0/en/From_Camera_Controlled_Video_to_Spatial_World_Models_Recomo_White_Paper_v2.0_EN.md",
    ROOT / "docs/spatial-world-model-white-paper/v2.0/zh-CN/从相机可控视频到空间世界模型_Recomo技术白皮书_v2.0_CN.md",
]

RESEARCH_REPORTS = [
    ROOT / "docs/camera-aware-video-generation/research-report/v1.0/en/camera_aware_video_generation_research_v1.0_EN.md",
    ROOT / "docs/camera-aware-video-generation/research-report/v1.0/zh-CN/相机感知视频生成与世界模型研究_v1.0_CN.md",
]

SLIDE_SOURCES = [
    ROOT / "docs/camera-aware-video-generation/research-report/v1.0/slides/camera_aware_video_generation_evolution_recommendations_v1.0_EN.md",
    ROOT / "docs/camera-aware-video-generation/research-report/v1.0/slides/相机感知视频生成演进与Recomo建议_v1.0_CN.md",
]


@dataclass
class Block:
    kind: str
    text: str = ""
    level: int = 0
    rows: list[list[str]] | None = None


def is_cjk(path: Path) -> bool:
    return "zh-CN" in path.parts or "_CN" in path.name or any("\u4e00" <= c <= "\u9fff" for c in path.name)


def strip_frontmatter(text: str) -> str:
    if text.startswith("---\n"):
        end = text.find("\n---\n", 4)
        if end >= 0:
            return text[end + 5 :]
    return text


def clean_inline(text: str) -> str:
    text = re.sub(r"!\[([^]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    return text.strip()


def parse_markdown(text: str) -> list[Block]:
    lines = strip_frontmatter(text).splitlines()
    blocks: list[Block] = []
    i = 0
    in_code = False
    code_lines: list[str] = []
    para: list[str] = []

    def flush_para() -> None:
        nonlocal para
        if para:
            joined = " ".join(x.strip() for x in para).strip()
            if joined:
                blocks.append(Block("paragraph", joined))
            para = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("```"):
            flush_para()
            if in_code:
                blocks.append(Block("code", "\n".join(code_lines)))
                code_lines = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(line)
            i += 1
            continue
        if not stripped:
            flush_para()
            i += 1
            continue
        if stripped == "---":
            flush_para()
            blocks.append(Block("rule"))
            i += 1
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            flush_para()
            blocks.append(Block("heading", clean_inline(heading.group(2)), len(heading.group(1))))
            i += 1
            continue
        if stripped.startswith("|") and i + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-{3,}", lines[i + 1]):
            flush_para()
            table_lines = [line]
            i += 1
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            rows = []
            for idx, raw in enumerate(table_lines):
                cells = [clean_inline(c) for c in raw.strip().strip("|").split("|")]
                if idx == 1 and all(re.fullmatch(r":?-{3,}:?", c.replace(" ", "")) for c in cells):
                    continue
                rows.append(cells)
            blocks.append(Block("table", rows=rows))
            continue
        bullet = re.match(r"^\s*[-*+]\s+(.*)$", line)
        number = re.match(r"^\s*\d+[.)]\s+(.*)$", line)
        if bullet or number:
            flush_para()
            blocks.append(Block("bullet" if bullet else "number", clean_inline((bullet or number).group(1))))
            i += 1
            continue
        if stripped.startswith(">"):
            flush_para()
            blocks.append(Block("quote", clean_inline(stripped.lstrip("> "))))
            i += 1
            continue
        if stripped.startswith("\\[") or stripped.startswith("\\("):
            flush_para()
            eq = [line]
            closing = "\\]" if stripped.startswith("\\[") else "\\)"
            while closing not in eq[-1] and i + 1 < len(lines):
                i += 1
                eq.append(lines[i])
            blocks.append(Block("equation", "\n".join(eq)))
            i += 1
            continue
        para.append(line)
        i += 1

    flush_para()
    if code_lines:
        blocks.append(Block("code", "\n".join(code_lines)))
    return blocks


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_run_font(run, font_name: str, size: float | None = None, bold: bool | None = None) -> None:
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold


def build_docx(source: Path, output: Path) -> None:
    cjk = is_cjk(source)
    body_font = "Noto Sans CJK SC" if cjk else "Aptos"
    heading_font = "Noto Sans CJK SC" if cjk else "Aptos Display"
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.72)
    sec.right_margin = Inches(0.72)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = body_font
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), body_font)
    normal.font.size = Pt(9.5 if not cjk else 10)
    normal.paragraph_format.space_after = Pt(5)
    normal.paragraph_format.line_spacing = 1.15

    for level in range(1, 5):
        sty = styles[f"Heading {level}"]
        sty.font.name = heading_font
        sty._element.rPr.rFonts.set(qn("w:eastAsia"), heading_font)
        sty.font.color.rgb = RGBColor(25, 45, 72)
        sty.font.size = Pt({1: 22, 2: 16, 3: 12.5, 4: 11}.get(level, 10.5))
        sty.font.bold = True
        sty.paragraph_format.space_before = Pt(10 if level <= 2 else 6)
        sty.paragraph_format.space_after = Pt(4)

    blocks = parse_markdown(source.read_text(encoding="utf-8"))
    title_seen = False
    for block in blocks:
        if block.kind == "heading":
            p = doc.add_paragraph(style=f"Heading {min(block.level, 4)}")
            r = p.add_run(block.text)
            set_run_font(r, heading_font, bold=True)
            if block.level == 1 and not title_seen:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.space_after = Pt(12)
                title_seen = True
        elif block.kind == "paragraph":
            p = doc.add_paragraph()
            r = p.add_run(clean_inline(block.text))
            set_run_font(r, body_font)
        elif block.kind in {"bullet", "number"}:
            style_name = "List Bullet" if block.kind == "bullet" else "List Number"
            p = doc.add_paragraph(style=style_name)
            r = p.add_run(block.text)
            set_run_font(r, body_font)
        elif block.kind == "quote":
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.28)
            p.paragraph_format.right_indent = Inches(0.15)
            r = p.add_run(block.text)
            set_run_font(r, body_font)
            r.italic = True
        elif block.kind in {"code", "equation"}:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.18)
            p.paragraph_format.space_before = Pt(3)
            p.paragraph_format.space_after = Pt(5)
            r = p.add_run(block.text)
            set_run_font(r, "Noto Sans Mono CJK SC" if cjk else "Consolas", 8.5)
        elif block.kind == "table" and block.rows:
            width = max(len(r) for r in block.rows)
            rows = [r + [""] * (width - len(r)) for r in block.rows]
            table = doc.add_table(rows=len(rows), cols=width)
            table.style = "Table Grid"
            for ri, row in enumerate(rows):
                for ci, value in enumerate(row):
                    cell = table.cell(ri, ci)
                    cell.text = ""
                    run = cell.paragraphs[0].add_run(value)
                    set_run_font(run, body_font, 8.5, bold=(ri == 0))
                    if ri == 0:
                        set_cell_shading(cell, "DCE6F1")
            doc.add_paragraph().paragraph_format.space_after = Pt(2)
        elif block.kind == "rule":
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(3)
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "4")
            bottom.set(qn("w:color"), "B7C7D9")
            pBdr.append(bottom)
            pPr.append(pBdr)

    core = doc.core_properties
    core.title = next((b.text for b in blocks if b.kind == "heading" and b.level == 1), source.stem)
    core.subject = "Recomo camera-aware video generation and spatial world-model research"
    core.author = "Recomo Research"
    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(output)


def find_font(cjk: bool) -> tuple[str, str]:
    candidates = []
    if cjk:
        candidates += [
            Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
            Path("/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf"),
            Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
        ]
    candidates += [
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf"),
    ]
    regular = next((p for p in candidates if p.exists()), None)
    bold_candidates = [
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf"),
    ]
    bold = next((p for p in bold_candidates if p.exists()), regular)
    if regular is None:
        return "Helvetica", "Helvetica-Bold"
    try:
        pdfmetrics.registerFont(TTFont("RecomoBody", str(regular), subfontIndex=0))
        pdfmetrics.registerFont(TTFont("RecomoBold", str(bold), subfontIndex=0))
        return "RecomoBody", "RecomoBold"
    except Exception:
        if Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf").exists():
            pdfmetrics.registerFont(TTFont("RecomoBody", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"))
            pdfmetrics.registerFont(TTFont("RecomoBold", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"))
            return "RecomoBody", "RecomoBold"
        return "Helvetica", "Helvetica-Bold"


def build_pdf(source: Path, output: Path) -> None:
    cjk = is_cjk(source)
    body_font, bold_font = find_font(cjk)
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "Body",
        parent=styles["BodyText"],
        fontName=body_font,
        fontSize=8.6 if not cjk else 9.1,
        leading=12.2 if not cjk else 13.8,
        textColor=colors.HexColor("#25364A"),
        spaceAfter=4,
    )
    bullets = ParagraphStyle("Bullet", parent=body, leftIndent=12, firstLineIndent=-7, bulletIndent=0)
    quote = ParagraphStyle(
        "Quote", parent=body, leftIndent=14, rightIndent=8, textColor=colors.HexColor("#405A73"), borderColor=colors.HexColor("#7B9DBD"), borderWidth=1, borderPadding=5
    )
    code = ParagraphStyle("Code", parent=body, fontName=body_font, fontSize=7.3, leading=9.6, leftIndent=8, backColor=colors.HexColor("#F1F5F8"), borderPadding=5)
    heads = {
        1: ParagraphStyle("H1", parent=body, fontName=bold_font, fontSize=22, leading=27, alignment=TA_CENTER, textColor=colors.HexColor("#142D48"), spaceBefore=8, spaceAfter=13),
        2: ParagraphStyle("H2", parent=body, fontName=bold_font, fontSize=15, leading=19, textColor=colors.HexColor("#173F67"), spaceBefore=11, spaceAfter=6),
        3: ParagraphStyle("H3", parent=body, fontName=bold_font, fontSize=11.5, leading=14, textColor=colors.HexColor("#285D82"), spaceBefore=7, spaceAfter=4),
        4: ParagraphStyle("H4", parent=body, fontName=bold_font, fontSize=9.7, leading=12, textColor=colors.HexColor("#285D82"), spaceBefore=5, spaceAfter=3),
    }

    def footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#B8C8D8"))
        canvas.line(18 * mm, 14 * mm, 192 * mm, 14 * mm)
        canvas.setFont(body_font, 7)
        canvas.setFillColor(colors.HexColor("#61758A"))
        canvas.drawString(18 * mm, 9.5 * mm, "Recomo Research")
        canvas.drawRightString(192 * mm, 9.5 * mm, str(doc.page))
        canvas.restoreState()

    doc = BaseDocTemplate(
        str(output), pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm, topMargin=17 * mm, bottomMargin=19 * mm,
        title=source.stem, author="Recomo Research"
    )
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")
    doc.addPageTemplates([PageTemplate(id="main", frames=[frame], onPage=footer)])

    story = []
    for block in parse_markdown(source.read_text(encoding="utf-8")):
        if block.kind == "heading":
            story.append(Paragraph(html.escape(block.text), heads[min(block.level, 4)]))
        elif block.kind == "paragraph":
            story.append(Paragraph(html.escape(clean_inline(block.text)), body))
        elif block.kind in {"bullet", "number"}:
            mark = "•" if block.kind == "bullet" else "–"
            story.append(Paragraph(f"{mark} {html.escape(block.text)}", bullets))
        elif block.kind == "quote":
            story.append(Paragraph(html.escape(block.text), quote))
        elif block.kind in {"code", "equation"}:
            story.append(Preformatted(block.text, code, maxLineLength=110))
        elif block.kind == "rule":
            story.append(Spacer(1, 4))
        elif block.kind == "table" and block.rows:
            width = max(len(r) for r in block.rows)
            rows = [r + [""] * (width - len(r)) for r in block.rows]
            data = [[Paragraph(html.escape(cell), ParagraphStyle("Cell", parent=body, fontSize=7.2, leading=9.2)) for cell in row] for row in rows]
            col_widths = [doc.width / width] * width
            table = Table(data, colWidths=col_widths, repeatRows=1, hAlign="LEFT")
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#DCE8F2")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#173F67")),
                ("FONTNAME", (0, 0), (-1, 0), bold_font),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#AABCCB")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]))
            story.extend([table, Spacer(1, 5)])
    output.parent.mkdir(parents=True, exist_ok=True)
    doc.build(story)


def split_slides(text: str) -> list[str]:
    text = strip_frontmatter(text)
    chunks = re.split(r"\n---\n", text)
    return [c.strip() for c in chunks if re.search(r"^#\s+", c, re.M)]


def parse_slide(chunk: str) -> tuple[str, list[str]]:
    lines = chunk.splitlines()
    title = ""
    body: list[str] = []
    in_code = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if not title and stripped.startswith("# "):
            title = clean_inline(stripped[2:])
            continue
        if not stripped:
            continue
        if stripped.startswith("## "):
            body.append(clean_inline(stripped[3:]))
        elif re.match(r"^[-*+]\s+", stripped):
            body.append("• " + clean_inline(re.sub(r"^[-*+]\s+", "", stripped)))
        elif re.match(r"^\d+[.)]\s+", stripped):
            body.append(clean_inline(stripped))
        else:
            body.append(clean_inline(stripped))
    return title, body


def build_pptx(source: Path, output: Path) -> None:
    cjk = is_cjk(source)
    font = "Noto Sans CJK SC" if cjk else "Aptos"
    prs = Presentation()
    prs.slide_width = PptxInches(13.333333)
    prs.slide_height = PptxInches(7.5)

    navy = PptxRGBColor(18, 42, 66)
    blue = PptxRGBColor(43, 104, 145)
    pale = PptxRGBColor(236, 243, 248)
    ink = PptxRGBColor(35, 54, 74)
    white = PptxRGBColor(255, 255, 255)
    muted = PptxRGBColor(99, 121, 140)

    chunks = split_slides(source.read_text(encoding="utf-8"))
    for index, chunk in enumerate(chunks):
        title, body = parse_slide(chunk)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = navy if index == 0 or index == len(chunks) - 1 else white

        if index not in {0, len(chunks) - 1}:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(0.22), prs.slide_height)
            band.fill.solid(); band.fill.fore_color.rgb = blue; band.line.fill.background()
            section = slide.shapes.add_textbox(PptxInches(11.7), PptxInches(0.27), PptxInches(1.0), PptxInches(0.3))
            p = section.text_frame.paragraphs[0]
            p.text = "RECOMO RESEARCH"
            p.font.name = font; p.font.size = PptxPt(7.5); p.font.bold = True; p.font.color.rgb = muted
            p.alignment = PP_ALIGN.RIGHT

        title_box = slide.shapes.add_textbox(PptxInches(0.65), PptxInches(0.58 if index else 1.55), PptxInches(12.0), PptxInches(1.15))
        tf = title_box.text_frame
        tf.clear(); tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = re.sub(r"^\d+\.\s*", "", title)
        p.font.name = font
        p.font.bold = True
        p.font.size = PptxPt(30 if index not in {0, len(chunks)-1} else 36)
        p.font.color.rgb = white if index in {0, len(chunks)-1} else navy

        if index == 0:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.67), PptxInches(3.0), PptxInches(2.0), PptxInches(0.08))
            accent.fill.solid(); accent.fill.fore_color.rgb = PptxRGBColor(88, 176, 204); accent.line.fill.background()

        body_box = slide.shapes.add_textbox(PptxInches(0.75), PptxInches(2.0 if index not in {0, len(chunks)-1} else 3.35), PptxInches(11.75), PptxInches(4.75 if index not in {0, len(chunks)-1} else 2.5))
        btf = body_box.text_frame
        btf.clear(); btf.word_wrap = True
        btf.margin_left = PptxInches(0.08); btf.margin_right = PptxInches(0.08)
        btf.margin_top = PptxInches(0.05); btf.margin_bottom = PptxInches(0.05)
        btf.vertical_anchor = MSO_ANCHOR.TOP
        max_lines = max(1, len(body))
        base_size = 20 if max_lines <= 8 else 17 if max_lines <= 13 else 14.5
        for j, item in enumerate(body):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            if item.startswith("• "):
                p.text = item[2:]
                p.level = 0
                p.text = "• " + p.text
            else:
                p.text = item
            p.font.name = font
            p.font.size = PptxPt(base_size)
            p.font.color.rgb = pale if index in {0, len(chunks)-1} else ink
            p.space_after = PptxPt(8 if max_lines <= 10 else 4)
            p.line_spacing = 1.08

        num = slide.shapes.add_textbox(PptxInches(12.35), PptxInches(7.08), PptxInches(0.45), PptxInches(0.22))
        p = num.text_frame.paragraphs[0]
        p.text = str(index + 1)
        p.font.name = font; p.font.size = PptxPt(8); p.font.color.rgb = pale if index in {0, len(chunks)-1} else muted
        p.alignment = PP_ALIGN.RIGHT

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    sources = WHITE_PAPERS + RESEARCH_REPORTS
    missing = [p for p in sources + SLIDE_SOURCES if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing canonical sources:\n" + "\n".join(str(p) for p in missing))

    for src in sources:
        build_docx(src, src.with_suffix(".docx"))
        build_pdf(src, src.with_suffix(".pdf"))
        print(f"built {src.with_suffix('.docx').relative_to(ROOT)}")
        print(f"built {src.with_suffix('.pdf').relative_to(ROOT)}")

    for src in SLIDE_SOURCES:
        build_pptx(src, src.with_suffix(".pptx"))
        print(f"built {src.with_suffix('.pptx').relative_to(ROOT)}")


if __name__ == "__main__":
    main()
